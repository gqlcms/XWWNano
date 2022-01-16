import  os
from pptx import Presentation
from pptx.util import Cm #Inches

# setting path

def Make_pptx( **kwargs):
    work_path = kwargs.get("work_path", None)
    pptxname = kwargs.get("pptxname", None)
    plotslist = kwargs.get("plotslist", None)
    Plot_Text = kwargs.get("Plot_Text", None)
    
    Number_of_PlotsIn_OneLine = kwargs.get("Number_of_PlotsIn_OneLine", (3,2)) # like (3,2)
    Plots_Size = kwargs.get("",6.5)
    Height_Left_ForText = kwargs.get("Height_Left_ForText",5)
    Left_Left_ForText = kwargs.get("Height_Left_ForText",2)
    Space_between_Plots = kwargs.get("Space_between_Plots",1)
    Space_between_Lines = kwargs.get("Space_between_Lines",1)
    Cone_Leftfor_Text = kwargs.get("Cone_Leftfor_Text",1)
    Input_Templete = kwargs.get("Input_Templete","/eos/user/q/qiguo/www/gKK/plots/templete/templete.pptx")
    
    Nplots_in_OneSlide = Number_of_PlotsIn_OneLine[0]*Number_of_PlotsIn_OneLine[1]

    if not os.path.isdir(work_path):
        os.mkdir(work_path)

    prs = Presentation(Input_Templete)

    slide_id_list = []
    for slide in prs.slides:
        slide_id_list.append(slide.slide_id)

    slide_1 = prs.slides.get(slide_id_list[0])

    slides_number = 1 ; count_plots = 0
    slide_workingon = slide_1

    

    for i in range(len(plotslist)):
        count_plots += 1
        if plotslist[i] == "split" :
            count_plots = 0
            slides_number += 1
            slide_id = slide_id_list[slides_number-1]
            # slide_workingon = prs.slides.add_slide(prs.slides.get(slide_id))
            slide_workingon = prs.slides.get(slide_id)
            continue

        

        img_name  = plotslist[i]
        print "adding ", img_name, "\n"
        
        # position of the plots
        if count_plots > Number_of_PlotsIn_OneLine[0]:
            position_x = Cm((count_plots - Number_of_PlotsIn_OneLine[0]-1)*(Plots_Size+Space_between_Plots) + Left_Left_ForText)
        else:
            position_x = Cm((count_plots-1)*(Plots_Size+Space_between_Plots) + Left_Left_ForText)

        position_y = Cm( (int( float(count_plots-1)/float(Number_of_PlotsIn_OneLine[0])) ) *(Plots_Size+Space_between_Lines) + Height_Left_ForText )

        if Plot_Text:
            position_yText = Cm( (int( float(count_plots-1)/float(Number_of_PlotsIn_OneLine[0])) ) *(Plots_Size+Space_between_Lines) + Height_Left_ForText - Cone_Leftfor_Text)
            txBox = slide_workingon.shapes.add_textbox(position_x, position_yText, Plots_Size, 4)
            tf = txBox.text_frame
            tf.text = Plot_Text[i]

        height = Cm(Plots_Size)
        width = Cm(Plots_Size)

        slide_workingon.shapes.add_picture(image_file=img_name,
                                left = position_x,
                                top = position_y,
                                height = height,
                                width = width,
                                )

    prs.save(work_path+pptxname)

    print "generate : ",work_path+pptxname

def get_file_list(file_path):
    dir_list = os.listdir(file_path)
    if not dir_list:
        return
    else:
        dir_list = sorted(dir_list,key=lambda x: os.path.getmtime(os.path.join(file_path, x)))
        return dir_list

def Wrapper_Plots_underOnefolder(InputPath, NPlots, selectplots = None, Comment = True):
    Allplots_0 = [ i for i in get_file_list(InputPath) if ".png" in i ]
    Allplots = []
    if selectplots:
        for j in selectplots:
            for i in Allplots_0:
                if j in i:
                    if i not in Allplots:
                        Allplots.append(i)
    else:
        Allplots = Allplots_0
    
    pptxname = InputPath.split("/")[-2]+".pptx"
    plotslist = [] ; Plot_Text = [];
    for i,j in enumerate(Allplots) :
        if (i+1)%NPlots == 0 :
            plotslist.append(InputPath+j)
            if Comment:
                Plot_Text.append(j.replace(".png","").replace("_"," "))
            else:
                Plot_Text.append(' ')
            plotslist.append('split')
            Plot_Text.append('split')
        else:
            plotslist.append(InputPath+j)
            if Comment:
                Plot_Text.append(j.replace(".png","").replace("_"," "))
            else:
                Plot_Text.append(' ')

    Make_pptx( pptxname = pptxname, plotslist = plotslist, work_path = InputPath, Plot_Text = Plot_Text)

    return Allplots



selectplots = None
Wrapper_Plots_underOnefolder("/eos/user/q/qiguo/www/VVV/plots/plots/PS/dataVsMC/", 6, selectplots = selectplots,  Comment = False)
